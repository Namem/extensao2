"""
Sprint Review PowerPoint — Ceres Diagnóstico — Sprint 1
Uses python-pptx to create the presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
import os
import lxml.etree as etree

# ── Color Palette ──────────────────────────────────────────────────────────────
C_FOREST_DARK   = RGBColor(0x1B, 0x3A, 0x0F)   # deep forest (dark bg)
C_FOREST        = RGBColor(0x2D, 0x50, 0x16)   # forest green
C_MOSS          = RGBColor(0x6B, 0x8F, 0x47)   # moss/olive
C_GOLD          = RGBColor(0xE8, 0xA0, 0x20)   # gold accent
C_CREAM         = RGBColor(0xF7, 0xF4, 0xEF)   # cream bg
C_WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT_DARK     = RGBColor(0x1A, 0x1A, 0x1A)
C_TEXT_MID      = RGBColor(0x4A, 0x4A, 0x4A)
C_SUCCESS       = RGBColor(0x2E, 0x7D, 0x32)
C_WARNING       = RGBColor(0xE6, 0x51, 0x00)

# ── Slide dimensions (16x9) ────────────────────────────────────────────────────
W = Inches(10)
H = Inches(5.625)

OUTPUT = r"C:\Users\Rachid\Desktop\NR\Semestre 2026_1\extensao\ceres-diagnostico\Pre_arquivos\SprintReview_1_CeresDiagnostico.pptx"


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ── Helper utilities ────────────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=14, bold=False, italic=False, color=None,
             align=PP_ALIGN.LEFT, v_align=None, font_face="Calibri",
             word_wrap=True, margin_left=0.08, margin_top=0.04,
             margin_right=0.08, margin_bottom=0.04):
    """Add a text box."""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if margin_left is not None:
        tf.margin_left   = Inches(margin_left)
        tf.margin_right  = Inches(margin_right)
        tf.margin_top    = Inches(margin_top)
        tf.margin_bottom = Inches(margin_bottom)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_face
    if color:
        run.font.color.rgb = color
    return txBox


def add_text_box_rich(slide, runs_list, x, y, w, h,
                      align=PP_ALIGN.LEFT, v_align=None,
                      margin_left=0.08, margin_top=0.04,
                      margin_right=0.08, margin_bottom=0.04,
                      word_wrap=True):
    """
    Add a text box with multiple paragraphs.
    runs_list: list of dicts:
      { 'text': str, 'size': int, 'bold': bool, 'color': RGBColor,
        'italic': bool, 'face': str, 'new_para': bool,
        'align': PP_ALIGN, 'space_after': int }
    """
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = Inches(margin_left)
    tf.margin_right  = Inches(margin_right)
    tf.margin_top    = Inches(margin_top)
    tf.margin_bottom = Inches(margin_bottom)

    first = True
    current_para = None

    for item in runs_list:
        if item.get('new_para', False) or first:
            if not first:
                current_para = tf.add_paragraph()
            else:
                current_para = tf.paragraphs[0]
                first = False
            a = item.get('align', align)
            current_para.alignment = a
            spa = item.get('space_after', 0)
            if spa:
                current_para.space_after = Pt(spa)
        run = current_para.add_run()
        run.text = item.get('text', '')
        run.font.size  = Pt(item.get('size', 14))
        run.font.bold  = item.get('bold', False)
        run.font.italic = item.get('italic', False)
        run.font.name  = item.get('face', 'Calibri')
        c = item.get('color', C_TEXT_DARK)
        if c:
            run.font.color.rgb = c

    return txBox


def section_tag(slide, label, x=0.45, y=0.22, w=2.8, h=0.28):
    """Green pill/rect with white label text."""
    r = add_rect(slide, x, y, w, h, C_FOREST)
    add_text(slide, label, x + 0.08, y + 0.02, w - 0.1, h - 0.04,
             font_size=9, bold=True, color=C_WHITE,
             align=PP_ALIGN.LEFT,
             margin_left=0.05, margin_top=0.02, margin_right=0.05, margin_bottom=0.02)
    return r


def gold_line(slide, x, y, w, thickness_pt=2.5):
    """Horizontal gold separator line."""
    line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(thickness_pt))
    line.fill.solid()
    line.fill.fore_color.rgb = C_GOLD
    line.line.fill.background()
    return line


def green_card(slide, x, y, w, h, title, body, title_size=12, body_size=10):
    """Forest green card with white text."""
    add_rect(slide, x, y, w, h, C_FOREST)
    # title
    add_text(slide, title, x + 0.12, y + 0.1, w - 0.24, 0.28,
             font_size=title_size, bold=True, color=C_WHITE,
             align=PP_ALIGN.LEFT)
    # body
    add_text(slide, body, x + 0.12, y + 0.38, w - 0.24, h - 0.5,
             font_size=body_size, bold=False, color=RGBColor(0xD4, 0xE6, 0xC3),
             align=PP_ALIGN.LEFT, word_wrap=True)
    return slide


def stat_box(slide, x, y, w, h, number, label):
    """Forest-green stat box with large number and small label."""
    add_rect(slide, x, y, w, h, C_FOREST)
    add_text(slide, number, x, y + 0.12, w, 0.55,
             font_size=34, bold=True, color=C_GOLD,
             align=PP_ALIGN.CENTER, margin_left=0, margin_right=0,
             margin_top=0, margin_bottom=0)
    add_text(slide, label, x + 0.08, y + 0.65, w - 0.16, 0.45,
             font_size=10, bold=False, color=C_WHITE,
             align=PP_ALIGN.CENTER, word_wrap=True)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide1_capa(prs):
    """Slide 1 — Capa (dark forest bg)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_FOREST_DARK

    # Subtle left accent strip
    add_rect(slide, 0, 0, 0.08, 5.625, C_FOREST)

    # Gold top bar
    gold_line(slide, 0, 0, 10, 4)

    # Gold horizontal separator
    gold_line(slide, 0.5, 2.35, 9.0, 2.5)

    # Main title
    add_text(slide, "CERES DIAGNÓSTICO",
             0.5, 0.7, 9.0, 1.1,
             font_size=52, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, font_face="Calibri",
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    # Subtitle
    add_text(slide, "Sprint Review — Sprint 1",
             0.5, 1.85, 9.0, 0.55,
             font_size=24, bold=False, color=C_GOLD,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    # Decorative corner rectangle (bottom right) — stays within footer area
    add_rect(slide, 7.5, 4.75, 2.5, 0.375, C_FOREST)

    # Bottom info
    add_text(slide,
             "Projeto de Sistemas Inteligentes  |  Eng. Computação",
             0.5, 2.75, 9.0, 0.42,
             font_size=14, bold=False, color=RGBColor(0xD4, 0xE6, 0xC3),
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    add_text(slide,
             "Sorriso-MT  |  Abril de 2026",
             0.5, 3.18, 9.0, 0.38,
             font_size=13, bold=False, color=C_MOSS,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    # Moss green leaf motif strip (bottom)
    add_rect(slide, 0, 5.1, 10, 0.525, C_FOREST)
    add_text(slide, "Bacharelado em Engenharia da Computação — IFMT Sorriso",
             0.5, 5.1, 9.0, 0.525,
             font_size=10, bold=False, color=C_MOSS,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0.08, margin_right=0, margin_bottom=0)

    return slide


def slide2_problema(prs):
    """Slide 2 — O Problema."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_CREAM

    # Top accent bar
    add_rect(slide, 0, 0, 10, 0.18, C_FOREST_DARK)

    section_tag(slide, "01 | O PROBLEMA", x=0.45, y=0.28, w=2.5, h=0.28)

    # Left column heading
    add_text(slide, "O Agricultor Familiar de Sorriso-MT",
             0.45, 0.72, 5.8, 0.55,
             font_size=20, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.LEFT,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    # Body paragraph
    add_text(slide,
             "Pragas e doenças no tomateiro causam perdas significativas "
             "de produção. O diagnóstico tardio ou incorreto impede o "
             "manejo adequado e onera o produtor.",
             0.45, 1.32, 5.7, 0.72,
             font_size=12, bold=False, color=C_TEXT_MID,
             align=PP_ALIGN.LEFT, word_wrap=True)

    # 3 bullet pain points
    bullets_data = [
        "Acesso restrito a agrônomos especializados no campo",
        "Sensibilidade ao tempo — diagnóstico tardio amplia perdas",
        "Custo elevado de laudos e insumos mal indicados",
    ]
    by = 2.1
    for b in bullets_data:
        # bullet circle
        add_rect(slide, 0.45, by + 0.06, 0.14, 0.14, C_GOLD)
        add_text(slide, b, 0.68, by, 5.4, 0.38,
                 font_size=12, bold=False, color=C_TEXT_DARK,
                 align=PP_ALIGN.LEFT, word_wrap=True,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)
        by += 0.44

    # Vertical divider
    add_rect(slide, 6.3, 0.65, 0.04, 4.6, C_MOSS)

    # Right column — stat boxes
    stats = [
        ("10",  "pragas/doenças\nmapeadas"),
        ("~3",  "perguntas até o\ndiagnóstico"),
        ("0",   "digitação\nnecessária"),
    ]
    sx = 6.55
    sy = 0.72
    for num, lbl in stats:
        stat_box(slide, sx, sy, 3.0, 1.18, num, lbl)
        sy += 1.3

    # Bottom bar
    add_rect(slide, 0, 5.25, 10, 0.375, C_FOREST_DARK)
    add_text(slide, "Ceres Diagnóstico — Sprint 1",
             0.45, 5.25, 9.1, 0.375,
             font_size=9, bold=False, color=C_MOSS,
             align=PP_ALIGN.RIGHT,
             margin_left=0, margin_top=0.06, margin_right=0.1, margin_bottom=0)

    return slide


def slide3_solucao(prs):
    """Slide 3 — A Solução."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_CREAM

    add_rect(slide, 0, 0, 10, 0.18, C_FOREST_DARK)
    section_tag(slide, "02 | A SOLUÇÃO", x=0.45, y=0.28, w=2.5, h=0.28)

    add_text(slide, "Sistema Especialista Determinístico",
             0.45, 0.72, 9.1, 0.55,
             font_size=22, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.LEFT,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    # 3 cards — reduced height so they don't crowd the bottom bar
    cards = [
        ("IA Determinística",
         "Árvore de decisão baseada em regras Embrapa — "
         "sem caixa-preta, auditável e previsível."),
        ("Mobile First",
         "Aplicativo Flutter com suporte offline via Drift — "
         "funciona sem internet na lavoura."),
        ("Open & Auditable",
         "Base de conhecimento em JSON, revisável e "
         "extensível por agrônomos."),
    ]
    cx = 0.45
    for title, body in cards:
        green_card(slide, cx, 1.52, 2.95, 2.9, title, body,
                   title_size=13, body_size=11)
        cx += 3.1

    # Gold accent at bottom of cards
    gold_line(slide, 0.45, 5.1, 9.1, 2.5)

    # Bottom bar
    add_rect(slide, 0, 5.25, 10, 0.375, C_FOREST_DARK)
    add_text(slide, "Ceres Diagnóstico — Sprint 1",
             0.45, 5.25, 9.1, 0.375,
             font_size=9, bold=False, color=C_MOSS,
             align=PP_ALIGN.RIGHT,
             margin_left=0, margin_top=0.06, margin_right=0.1, margin_bottom=0)

    return slide


def slide4_planejado(prs):
    """Slide 4 — Sprint 1: O que foi Planejado."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_CREAM

    add_rect(slide, 0, 0, 10, 0.18, C_FOREST_DARK)
    section_tag(slide, "03 | SPRINT 1 — PLANEJADO", x=0.45, y=0.28, w=3.5, h=0.28)

    add_text(slide, "6 Objetivos Definidos no Backlog",
             0.45, 0.72, 9.1, 0.52,
             font_size=22, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.LEFT,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    items = [
        ("1", "Setup do Ambiente",         "Docker, Django, Flutter, Git"),
        ("2", "Modelagem de Dados",        "Pergunta / Opcao / Diagnostico"),
        ("3", "Segurança JWT",             "SimpleJWT + endpoints de token"),
        ("4", "Motor de Inferência",       "GET iniciar/ + POST responder/"),
        ("5", "Testes Unitários",          "APITestCase — happy paths"),
        ("6", "Engenharia do Conhecimento","10 doenças Embrapa + JSON tree"),
    ]

    # 2-column × 3-row grid
    col_xs = [0.45, 5.1]
    row_ys = [1.45, 2.65, 3.85]

    for i, (num, title, subtitle) in enumerate(items):
        col = i % 2
        row = i // 2
        cx = col_xs[col]
        cy = row_ys[row]
        bw = 4.45

        # Number badge
        add_rect(slide, cx, cy + 0.08, 0.42, 0.42, C_GOLD)
        add_text(slide, num, cx, cy + 0.08, 0.42, 0.42,
                 font_size=16, bold=True, color=C_FOREST_DARK,
                 align=PP_ALIGN.CENTER,
                 margin_left=0, margin_top=0.02, margin_right=0, margin_bottom=0)

        # Card bg
        add_rect(slide, cx + 0.5, cy, bw - 0.5, 0.92, C_WHITE,
                 line_color=C_MOSS, line_width=0.8)

        # Title
        add_text(slide, title,
                 cx + 0.62, cy + 0.04, bw - 0.74, 0.36,
                 font_size=13, bold=True, color=C_FOREST_DARK,
                 align=PP_ALIGN.LEFT,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

        # Subtitle
        add_text(slide, subtitle,
                 cx + 0.62, cy + 0.40, bw - 0.74, 0.45,
                 font_size=10, bold=False, color=C_TEXT_MID,
                 align=PP_ALIGN.LEFT,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    # Bottom bar
    add_rect(slide, 0, 5.25, 10, 0.375, C_FOREST_DARK)
    add_text(slide, "Ceres Diagnóstico — Sprint 1",
             0.45, 5.25, 9.1, 0.375,
             font_size=9, bold=False, color=C_MOSS,
             align=PP_ALIGN.RIGHT,
             margin_left=0, margin_top=0.06, margin_right=0.1, margin_bottom=0)

    return slide


def slide5_entregue(prs):
    """Slide 5 — Sprint 1: O que foi Entregue."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_CREAM

    add_rect(slide, 0, 0, 10, 0.18, C_FOREST_DARK)
    section_tag(slide, "04 | SPRINT 1 — ENTREGUE", x=0.45, y=0.28, w=3.4, h=0.28)

    add_text(slide, "Conformidade com o Backlog",
             0.45, 0.72, 9.1, 0.52,
             font_size=22, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.LEFT,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    rows_data = [
        ("Setup do Ambiente",          "COMPLETO", "Docker + Django + Flutter + Git configurados"),
        ("Modelagem de Dados",         "COMPLETO", "Models Pergunta/Opcao/Diagnostico + migrações"),
        ("Segurança JWT",              "COMPLETO", "SimpleJWT configurado; /api/auth/token/ registrado"),
        ("Motor de Inferência",        "COMPLETO", "GET iniciar/ + POST responder/ funcionais"),
        ("Testes Unitários",           "COMPLETO", "3 APITestCase cobrindo iniciar + responder"),
        ("Engenharia do Conhecimento", "COMPLETO", "10 doenças Embrapa; JSON tree + script import"),
    ]

    header_y = 1.32
    row_h = 0.52
    col_widths = [3.0, 1.5, 4.6]
    col_xs = [0.45, 3.55, 5.1]

    # Header row
    headers = ["Objetivo", "Status", "Evidência"]
    for ci, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_widths)):
        add_rect(slide, cx, header_y, cw - 0.06, 0.38, C_FOREST_DARK)
        add_text(slide, hdr, cx + 0.08, header_y, cw - 0.14, 0.38,
                 font_size=11, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT,
                 margin_left=0.06, margin_top=0.05, margin_right=0, margin_bottom=0)

    # Data rows
    row_y = header_y + 0.42
    for ri, (goal, status, evidence) in enumerate(rows_data):
        bg = C_WHITE if ri % 2 == 0 else RGBColor(0xEE, 0xF4, 0xE8)

        for ci, (text, cx, cw) in enumerate(zip([goal, status, evidence], col_xs, col_widths)):
            add_rect(slide, cx, row_y, cw - 0.06, row_h - 0.06, bg,
                     line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=0.5)

        # Goal
        add_text(slide, goal, col_xs[0] + 0.08, row_y + 0.04,
                 col_widths[0] - 0.2, row_h - 0.12,
                 font_size=10, bold=False, color=C_TEXT_DARK,
                 align=PP_ALIGN.LEFT,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

        # Status badge (green)
        add_rect(slide, col_xs[1] + 0.1, row_y + 0.1, 1.25, 0.28, C_SUCCESS)
        add_text(slide, "✔ " + status, col_xs[1] + 0.1, row_y + 0.1, 1.25, 0.28,
                 font_size=9, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER,
                 margin_left=0, margin_top=0.03, margin_right=0, margin_bottom=0)

        # Evidence
        add_text(slide, evidence, col_xs[2] + 0.08, row_y + 0.04,
                 col_widths[2] - 0.2, row_h - 0.12,
                 font_size=9, bold=False, color=C_TEXT_MID,
                 align=PP_ALIGN.LEFT, word_wrap=True,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

        row_y += row_h

    # Summary note — positioned just above the bottom bar with breathing room
    add_rect(slide, 0.45, 4.96, 4.4, 0.24, C_SUCCESS)
    add_text(slide, "  6 de 6 objetivos entregues",
             0.45, 4.96, 4.4, 0.24,
             font_size=10, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0.03, margin_right=0, margin_bottom=0)

    # Bottom bar
    add_rect(slide, 0, 5.25, 10, 0.375, C_FOREST_DARK)
    add_text(slide, "Ceres Diagnóstico — Sprint 1",
             0.45, 5.25, 9.1, 0.375,
             font_size=9, bold=False, color=C_MOSS,
             align=PP_ALIGN.RIGHT,
             margin_left=0, margin_top=0.06, margin_right=0.1, margin_bottom=0)

    return slide


def slide6_motor(prs):
    """Slide 6 — Motor de Inferência."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_CREAM

    add_rect(slide, 0, 0, 10, 0.18, C_FOREST_DARK)
    section_tag(slide, "05 | MOTOR DE INFERÊNCIA", x=0.45, y=0.28, w=3.5, h=0.28)

    add_text(slide, "Arquitetura Técnica",
             0.45, 0.72, 9.1, 0.52,
             font_size=22, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.LEFT,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    # Left: architecture layers
    layers = [
        (C_MOSS,         "Flutter (Dart)",      "UI — Perguntas + Resultado"),
        (C_FOREST,       "Django + DRF",        "Motor de Inferência + JWT"),
        (C_FOREST_DARK,  "PostgreSQL",           "Árvore de decisão persistida"),
    ]
    ly = 1.45
    for fill, title, sub in layers:
        add_rect(slide, 0.45, ly, 4.3, 0.82, fill)
        add_text(slide, title, 0.6, ly + 0.04, 4.0, 0.36,
                 font_size=13, bold=True, color=C_WHITE,
                 align=PP_ALIGN.LEFT,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)
        add_text(slide, sub, 0.6, ly + 0.38, 4.0, 0.38,
                 font_size=10, bold=False, color=RGBColor(0xCC, 0xDD, 0xBB),
                 align=PP_ALIGN.LEFT,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)
        ly += 0.92

        # Arrow down (except last)
        if fill != C_FOREST_DARK:
            add_rect(slide, 2.25, ly, 0.3, 0.22, C_GOLD)
            add_text(slide, "▼", 2.25, ly, 0.3, 0.22,
                     font_size=10, bold=True, color=C_FOREST_DARK,
                     align=PP_ALIGN.CENTER,
                     margin_left=0, margin_top=0.01, margin_right=0, margin_bottom=0)
            ly += 0.28

    # Right: API endpoints — compact spacing so JSON box fits above bottom bar
    add_rect(slide, 5.0, 1.35, 4.55, 0.32, C_FOREST_DARK)
    add_text(slide, "API CONTRACT", 5.0, 1.35, 4.55, 0.32,
             font_size=11, bold=True, color=C_GOLD,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0.04, margin_right=0, margin_bottom=0)

    endpoints = [
        ("GET",  "/api/diagnostico/iniciar/",  "Retorna a pergunta raiz"),
        ("POST", "/api/diagnostico/responder/","{'opcao_id': N} → pergunta ou diagnóstico"),
        ("POST", "/api/auth/token/",           "{'username','password'} → JWT pair"),
        ("POST", "/api/auth/token/refresh/",   "{'refresh'} → novo access token"),
    ]
    method_colors = {
        "GET":  RGBColor(0x2E, 0x7D, 0x32),
        "POST": RGBColor(0x01, 0x57, 0x9B),
    }
    ey = 1.74
    for method, path, desc in endpoints:
        add_rect(slide, 5.0, ey, 0.6, 0.25, method_colors[method])
        add_text(slide, method, 5.0, ey, 0.6, 0.25,
                 font_size=8, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER,
                 margin_left=0, margin_top=0.03, margin_right=0, margin_bottom=0)
        add_rect(slide, 5.62, ey, 3.93, 0.25, RGBColor(0x2A, 0x2A, 0x2A))
        add_text(slide, path, 5.64, ey, 3.8, 0.25,
                 font_size=9, bold=False, color=RGBColor(0xA8, 0xE6, 0xA8),
                 align=PP_ALIGN.LEFT, font_face="Consolas",
                 margin_left=0.04, margin_top=0.03, margin_right=0, margin_bottom=0)
        add_text(slide, desc, 5.0, ey + 0.27, 4.55, 0.25,
                 font_size=9, bold=False, color=C_TEXT_MID,
                 align=PP_ALIGN.LEFT, word_wrap=True,
                 margin_left=0.04, margin_top=0, margin_right=0, margin_bottom=0)
        ey += 0.62

    # Response envelope note — fits above bottom bar (ey should be ~4.22 max)
    add_rect(slide, 5.0, ey + 0.04, 4.55, 0.58, RGBColor(0x2A, 0x2A, 0x2A))
    add_text(slide,
             '{"tipo": "pergunta"|"diagnostico",\n "dados": { id, texto, opcoes/nome/recomendacao }}',
             5.04, ey + 0.06, 4.47, 0.54,
             font_size=8, bold=False, color=RGBColor(0xA8, 0xE6, 0xA8),
             align=PP_ALIGN.LEFT, font_face="Consolas", word_wrap=True,
             margin_left=0.04, margin_top=0.04, margin_right=0.04, margin_bottom=0)

    # Bottom bar
    add_rect(slide, 0, 5.25, 10, 0.375, C_FOREST_DARK)
    add_text(slide, "Ceres Diagnóstico — Sprint 1",
             0.45, 5.25, 9.1, 0.375,
             font_size=9, bold=False, color=C_MOSS,
             align=PP_ALIGN.RIGHT,
             margin_left=0, margin_top=0.06, margin_right=0.1, margin_bottom=0)

    return slide


def slide7_conhecimento(prs):
    """Slide 7 — Base de Conhecimento."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_CREAM

    add_rect(slide, 0, 0, 10, 0.18, C_FOREST_DARK)
    section_tag(slide, "06 | BASE DE CONHECIMENTO", x=0.45, y=0.28, w=3.6, h=0.28)

    add_text(slide, "10 Pragas & Doenças — Embrapa",
             0.45, 0.72, 9.1, 0.52,
             font_size=22, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.LEFT,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    # Left: tree structure
    tree_lines = [
        ("P1: Onde está o sintoma?", 0.45, 1.42, C_FOREST_DARK, True, 11),
        ("├─ Folhas →", 0.55, 1.78, C_MOSS, False, 10),
        ("│    P2: Tipo de mancha? →", 0.7, 2.05, C_TEXT_MID, False, 10),
        ("│    ├─ P3: Com halo amarelo?", 0.85, 2.30, C_TEXT_MID, False, 9),
        ("│    │    ├─ Sim → D1 Requeima", 1.0, 2.52, C_SUCCESS, False, 9),
        ("│    │    └─ Não → D2 Septoriose", 1.0, 2.72, C_SUCCESS, False, 9),
        ("│    └─ P6: Pequenas pontuações?", 0.85, 2.92, C_TEXT_MID, False, 9),
        ("│         └─ D3 Pinta-Preta", 1.0, 3.12, C_SUCCESS, False, 9),
        ("├─ Frutos → P8 → D9 / D10", 0.55, 3.38, C_MOSS, False, 10),
        ("└─ Planta toda → P9 → D4-D8", 0.55, 3.62, C_MOSS, False, 10),
    ]

    for text, tx, ty, col, bold, size in tree_lines:
        add_text(slide, text, tx, ty, 4.2, 0.3,
                 font_size=size, bold=bold, color=col,
                 align=PP_ALIGN.LEFT, font_face="Consolas",
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    # Divider
    add_rect(slide, 4.85, 1.35, 0.04, 3.7, C_MOSS)

    # Right: 10 diagnoses in 2 columns of 5
    diseases = [
        "D1  Requeima",
        "D2  Septoriose",
        "D3  Pinta-Preta",
        "D4  Traça-do-Tomateiro",
        "D5  Mosca-Branca",
        "D6  Tripes/Vira-Cabeça",
        "D7  Ácaro-do-Bronzeamento",
        "D8  Murcha-Bacteriana",
        "D9  Mancha-Bacteriana",
        "D10 Broca-Pequena-do-Fruto",
    ]
    d_col_xs = [5.05, 7.55]
    dy = 1.38
    for i, d in enumerate(diseases):
        col = i % 2
        row = i // 2
        dx = d_col_xs[col]
        dy_pos = 1.38 + row * 0.72

        add_rect(slide, dx, dy_pos, 2.3, 0.58,
                 C_FOREST if i % 2 == 0 else C_FOREST_DARK)
        add_text(slide, d, dx + 0.1, dy_pos + 0.06, 2.1, 0.46,
                 font_size=10, bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, font_face="Calibri",
                 margin_left=0.05, margin_top=0.04, margin_right=0, margin_bottom=0)

    # Bottom bar
    add_rect(slide, 0, 5.25, 10, 0.375, C_FOREST_DARK)
    add_text(slide, "Ceres Diagnóstico — Sprint 1",
             0.45, 5.25, 9.1, 0.375,
             font_size=9, bold=False, color=C_MOSS,
             align=PP_ALIGN.RIGHT,
             margin_left=0, margin_top=0.06, margin_right=0.1, margin_bottom=0)

    return slide


def slide8_demo(prs):
    """Slide 8 — Demo ao Vivo (dark bg)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_FOREST_DARK

    # Top accent
    add_rect(slide, 0, 0, 10, 0.18, C_GOLD)

    # Section tag on dark bg — use gold bg for high contrast
    add_rect(slide, 0.45, 0.28, 2.8, 0.28, C_GOLD)
    add_text(slide, "07 | DEMONSTRAÇÃO", 0.53, 0.30, 2.64, 0.24,
             font_size=9, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.LEFT,
             margin_left=0.05, margin_top=0.02, margin_right=0.05, margin_bottom=0.02)

    add_text(slide, "Motor de Inferência — Demo ao Vivo",
             0.45, 0.72, 9.1, 0.55,
             font_size=22, bold=True, color=C_WHITE,
             align=PP_ALIGN.LEFT,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    gold_line(slide, 0.45, 1.35, 9.1, 2)

    steps = [
        ("1", "GET /api/diagnostico/iniciar/",
         "Retorna pergunta raiz: «Onde está o sintoma?» com opções Folhas / Frutos / Planta toda"),
        ("2", "POST /api/diagnostico/responder/  { \"opcao_id\": 1 }",
         "Seleciona ramo Folhas → retorna P2: «Que tipo de mancha você observa?»"),
        ("3", "POST /api/diagnostico/responder/  { \"opcao_id\": <manchas com halo> }",
         "Avança para P3: «A mancha tem halo amarelo ao redor?» — responde Sim"),
        ("4", "POST /api/diagnostico/responder/  { \"opcao_id\": <sim> }",
         "Tipo: «diagnostico» → D1 Requeima + recomendação Embrapa completa"),
    ]

    sy = 1.55
    for num, cmd, desc in steps:
        # Number bubble
        add_rect(slide, 0.45, sy, 0.42, 0.42, C_GOLD)
        add_text(slide, num, 0.45, sy, 0.42, 0.42,
                 font_size=16, bold=True, color=C_FOREST_DARK,
                 align=PP_ALIGN.CENTER,
                 margin_left=0, margin_top=0.04, margin_right=0, margin_bottom=0)

        # Command (code style)
        add_rect(slide, 0.95, sy, 8.6, 0.3, RGBColor(0x0A, 0x1A, 0x05))
        add_text(slide, cmd, 0.98, sy, 8.54, 0.3,
                 font_size=9, bold=False, color=RGBColor(0xA8, 0xE6, 0xA8),
                 align=PP_ALIGN.LEFT, font_face="Consolas",
                 margin_left=0.05, margin_top=0.04, margin_right=0.05, margin_bottom=0)

        # Description
        add_text(slide, desc, 0.95, sy + 0.32, 8.6, 0.38,
                 font_size=10, bold=False, color=RGBColor(0xD4, 0xE6, 0xC3),
                 align=PP_ALIGN.LEFT, word_wrap=True,
                 margin_left=0.06, margin_top=0, margin_right=0.06, margin_bottom=0)

        sy += 0.86

    # Bottom note — gold accent bg for legibility on dark slide
    add_rect(slide, 0.45, 5.0, 9.1, 0.22, C_GOLD)
    add_text(slide, "Ferramentas: Postman + Docker + Django runserver  (localhost:8000)",
             0.5, 5.0, 9.0, 0.22,
             font_size=9, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0.03, margin_right=0, margin_bottom=0)

    # Bottom bar
    add_rect(slide, 0, 5.25, 10, 0.375, C_FOREST)
    add_text(slide, "Ceres Diagnóstico — Sprint 1",
             0.45, 5.25, 9.1, 0.375,
             font_size=9, bold=False, color=C_GOLD,
             align=PP_ALIGN.RIGHT,
             margin_left=0, margin_top=0.06, margin_right=0.1, margin_bottom=0)

    return slide


def slide9_nao_concluido(prs):
    """Slide 9 — O que Não Foi Concluído."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_CREAM

    add_rect(slide, 0, 0, 10, 0.18, C_FOREST_DARK)
    section_tag(slide, "08 | NÃO CONCLUÍDO", x=0.45, y=0.28, w=2.8, h=0.28)

    add_text(slide, "Transparência sobre o Incremento",
             0.45, 0.72, 9.1, 0.52,
             font_size=22, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.LEFT,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    items = [
        ("Flutter — Sem integração de API",
         "Scaffolding apenas; integração de API prevista para Sprint 2.",
         "Previsto para Sprint 2"),
        ("Testes — Cobertura parcial",
         "3 happy-path tests implementados; edge cases e cenários de erro "
         "planejados para Sprint 2.",
         "Edge cases — Sprint 2"),
    ]

    iy = 1.45
    for title, body, sprint2 in items:
        # Warning color left bar
        add_rect(slide, 0.45, iy, 0.12, 1.2, C_WARNING)

        # Card bg
        add_rect(slide, 0.6, iy, 8.95, 1.2,
                 RGBColor(0xFF, 0xF3, 0xE0),
                 line_color=C_WARNING, line_width=0.8)

        # Title
        add_text(slide, title, 0.75, iy + 0.06, 6.8, 0.36,
                 font_size=13, bold=True, color=C_WARNING,
                 align=PP_ALIGN.LEFT,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

        # Sprint 2 tag
        add_rect(slide, 7.9, iy + 0.06, 1.5, 0.28, C_FOREST)
        add_text(slide, sprint2, 7.9, iy + 0.06, 1.5, 0.28,
                 font_size=8, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER,
                 margin_left=0, margin_top=0.04, margin_right=0, margin_bottom=0)

        # Body
        add_text(slide, body, 0.75, iy + 0.44, 8.6, 0.68,
                 font_size=11, bold=False, color=C_TEXT_MID,
                 align=PP_ALIGN.LEFT, word_wrap=True,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

        iy += 1.38

    # Green success box
    add_rect(slide, 0.45, 4.25, 9.1, 0.75, RGBColor(0xE8, 0xF5, 0xE9),
             line_color=C_SUCCESS, line_width=1.2)
    add_rect(slide, 0.45, 4.25, 0.12, 0.75, C_SUCCESS)
    add_text(slide,
             "2 itens corrigidos antes desta Review: "
             "renomeação do arquivo de testes (teste.py → tests.py) "
             "+ registro das rotas JWT em ceres_core/urls.py.",
             0.65, 4.28, 8.7, 0.68,
             font_size=11, bold=False, color=C_SUCCESS,
             align=PP_ALIGN.LEFT, word_wrap=True,
             margin_left=0.05, margin_top=0, margin_right=0, margin_bottom=0)

    # Bottom bar
    add_rect(slide, 0, 5.25, 10, 0.375, C_FOREST_DARK)
    add_text(slide, "Ceres Diagnóstico — Sprint 1",
             0.45, 5.25, 9.1, 0.375,
             font_size=9, bold=False, color=C_MOSS,
             align=PP_ALIGN.RIGHT,
             margin_left=0, margin_top=0.06, margin_right=0.1, margin_bottom=0)

    return slide


def slide10_sprint2(prs):
    """Slide 10 — Próximos Passos — Sprint 2."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_CREAM

    add_rect(slide, 0, 0, 10, 0.18, C_FOREST_DARK)
    section_tag(slide, "09 | SPRINT 2", x=0.45, y=0.28, w=2.0, h=0.28)

    add_text(slide, "Próximos Passos",
             0.45, 0.72, 9.1, 0.52,
             font_size=22, bold=True, color=C_FOREST_DARK,
             align=PP_ALIGN.LEFT,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    cards = [
        ("01", "Telas Flutter",
         "Fluxo completo de diagnóstico: cards de perguntas com opções e tela de resultado com recomendação."),
        ("02", "Integração API",
         "Dart models gerados + cliente HTTP/Dio para consumir os endpoints do motor de inferência."),
        ("03", "Autenticação JWT",
         "Fluxo de login no Flutter com armazenamento seguro do token (flutter_secure_storage)."),
        ("04", "Offline — Drift",
         "Schema Drift para cache local da sessão de diagnóstico — funciona sem conectividade."),
    ]

    cx = 0.45
    for num, title, body in cards:
        # Card background
        add_rect(slide, cx, 1.48, 2.18, 3.5, C_WHITE,
                 line_color=C_MOSS, line_width=0.8)

        # Gold number top bar
        add_rect(slide, cx, 1.48, 2.18, 0.45, C_GOLD)
        add_text(slide, num, cx, 1.48, 2.18, 0.45,
                 font_size=20, bold=True, color=C_FOREST_DARK,
                 align=PP_ALIGN.CENTER,
                 margin_left=0, margin_top=0.04, margin_right=0, margin_bottom=0)

        # Moss accent left bar
        add_rect(slide, cx, 1.93, 0.1, 3.05, C_MOSS)

        # Title
        add_text(slide, title, cx + 0.18, 1.98, 1.9, 0.44,
                 font_size=12, bold=True, color=C_FOREST_DARK,
                 align=PP_ALIGN.LEFT, word_wrap=True,
                 margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

        # Body
        add_text(slide, body, cx + 0.18, 2.44, 1.9, 1.95,
                 font_size=10, bold=False, color=C_TEXT_MID,
                 align=PP_ALIGN.LEFT, word_wrap=True,
                 margin_left=0, margin_top=0, margin_right=0.05, margin_bottom=0)

        cx += 2.32

    # Bottom bar
    add_rect(slide, 0, 5.25, 10, 0.375, C_FOREST_DARK)
    add_text(slide, "Ceres Diagnóstico — Sprint 1",
             0.45, 5.25, 9.1, 0.375,
             font_size=9, bold=False, color=C_MOSS,
             align=PP_ALIGN.RIGHT,
             margin_left=0, margin_top=0.06, margin_right=0.1, margin_bottom=0)

    return slide


def slide11_encerramento(prs):
    """Slide 11 — Encerramento & Feedback (dark bg)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_FOREST_DARK

    add_rect(slide, 0, 0, 10, 0.18, C_GOLD)

    # Decorative side strip
    add_rect(slide, 0, 0, 0.18, 5.625, C_FOREST)

    add_text(slide, "Obrigado!",
             0.5, 0.4, 9.0, 0.9,
             font_size=48, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    add_text(slide, "Ceres Diagnóstico  |  Sprint 1 — Concluída",
             0.5, 1.35, 9.0, 0.44,
             font_size=18, bold=False, color=C_GOLD,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    gold_line(slide, 0.5, 1.88, 9.0, 2)

    # Feedback heading
    add_text(slide, "Perguntas para Validação",
             0.5, 1.94, 9.0, 0.38,
             font_size=14, bold=True, color=C_MOSS,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0, margin_right=0, margin_bottom=0)

    questions = [
        "A lógica da árvore de decisão atende ao problema do agricultor?",
        "A estrutura da API está pronta para ser consumida pelo Flutter?",
        "Há casos de uso ou sintomas não cobertos pelos 10 diagnósticos atuais?",
    ]

    # Spread questions evenly through remaining vertical space (2.38 to 4.9)
    qy = 2.38
    q_spacing = 0.82
    for i, q in enumerate(questions):
        # Full-width card bg for each question
        add_rect(slide, 0.5, qy, 9.0, 0.6, C_FOREST)

        # Question number in gold — left
        add_rect(slide, 0.5, qy, 0.6, 0.6, C_GOLD)
        add_text(slide, str(i + 1), 0.5, qy, 0.6, 0.6,
                 font_size=16, bold=True, color=C_FOREST_DARK,
                 align=PP_ALIGN.CENTER,
                 margin_left=0, margin_top=0.08, margin_right=0, margin_bottom=0)

        add_text(slide, q, 1.18, qy + 0.06, 8.24, 0.5,
                 font_size=12, bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, word_wrap=True,
                 margin_left=0.08, margin_top=0.04, margin_right=0.08, margin_bottom=0)
        qy += q_spacing

    # Bottom document reference
    add_rect(slide, 0, 5.1, 10, 0.525, C_FOREST)
    add_text(slide,
             "Roteiro de Validação: guia_sprint_review.pdf  —  Sorriso-MT | Abril de 2026",
             0.5, 5.1, 9.0, 0.525,
             font_size=10, bold=False, color=C_GOLD,
             align=PP_ALIGN.CENTER,
             margin_left=0, margin_top=0.1, margin_right=0, margin_bottom=0)

    return slide


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slides...")
    slide1_capa(prs)
    print("  [1/11] Capa — done")

    slide2_problema(prs)
    print("  [2/11] O Problema — done")

    slide3_solucao(prs)
    print("  [3/11] A Solução — done")

    slide4_planejado(prs)
    print("  [4/11] Sprint 1 Planejado — done")

    slide5_entregue(prs)
    print("  [5/11] Sprint 1 Entregue — done")

    slide6_motor(prs)
    print("  [6/11] Motor de Inferência — done")

    slide7_conhecimento(prs)
    print("  [7/11] Base de Conhecimento — done")

    slide8_demo(prs)
    print("  [8/11] Demo ao Vivo — done")

    slide9_nao_concluido(prs)
    print("  [9/11] Não Concluído — done")

    slide10_sprint2(prs)
    print("  [10/11] Sprint 2 — done")

    slide11_encerramento(prs)
    print("  [11/11] Encerramento — done")

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
