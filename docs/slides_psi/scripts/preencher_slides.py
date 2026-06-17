"""Preenche Ceres_Diagnostico_-_Defesa_PSI.pptx com imagens e vídeo."""

import os
import zipfile
import shutil
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

DIR = os.path.dirname(os.path.abspath(__file__))
INPUT  = os.path.join(DIR, "Ceres_Diagnostico_-_Defesa_PSI.pptx")
OUTPUT = os.path.join(DIR, "Ceres_Diagnostico_-_Defesa_PSI_preenchido.pptx")
VIDEO  = os.path.join(DIR, "demo_ceres_1.5x.mp4")

LEAF_IMGS = [
    os.path.join(DIR, "cls_01_requeima.jpg"),
    os.path.join(DIR, "cls_02_septoriose.jpg"),
    os.path.join(DIR, "cls_03_pinta_preta.jpg"),
    os.path.join(DIR, "cls_04_mancha_alvo.jpg"),
    os.path.join(DIR, "cls_05_mofo_foliar.jpg"),
    os.path.join(DIR, "cls_06_vira_cabeca.jpg"),
    os.path.join(DIR, "cls_07_mosaico.jpg"),
    os.path.join(DIR, "cls_08_acaro.jpg"),
    os.path.join(DIR, "cls_09_mancha_bact.jpg"),
    os.path.join(DIR, "cls_10_saudavel.jpg"),
]
APP_IOT  = os.path.join(DIR, "app_iot.jpg")
HARDWARE = os.path.join(DIR, "hardware_setup.jpg")
MATRIZ   = os.path.join(DIR, "matriz_confusao_int8.png")

# ── 1. Abrir PPTX com python-pptx ───────────────────────────────────────────
prs = Presentation(INPUT)

# ── 2. SLIDE 7 — 10 fotos de folha ──────────────────────────────────────────
slide7 = prs.slides[6]

placeholders = []
for shape in slide7.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip().upper()
    if "IMAGEM" in text and "FOLHA" in text:
        aspect = shape.width / shape.height if shape.height else 999
        if 0.4 < aspect < 4.0:          # exclui o banner largo do subtítulo
            placeholders.append(shape)

# ordenar por posição (linha e depois coluna)
placeholders.sort(key=lambda s: (round(s.top / 100000), s.left))
print(f"Slide 7 — {len(placeholders)} placeholders de imagem encontrados")

for idx, (shape, img_path) in enumerate(zip(placeholders, LEAF_IMGS)):
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    slide7.shapes.add_picture(img_path, left, top, width, height)
    print(f"  [{idx+1}] {os.path.basename(img_path)} OK ({left//914400:.2f}\", {top//914400:.2f}\")")

# ── 2b. SLIDE 10 — Matriz de Confusão ────────────────────────────────────────
slide10 = prs.slides[9]

# Shape 6 = retângulo vazio (placeholder), Shape 7 = texto "MATRIZ 10x10"
# Remover ambos e inserir a imagem no lugar
matriz_shapes = []
for shape in slide10.shapes:
    if shape.has_text_frame:
        txt = shape.text_frame.text.strip().upper()
        if "MATRIZ" in txt and "HEATMAP" in txt:
            matriz_shapes.append(shape)
    elif (not shape.has_text_frame
          and shape.left == 762000
          and shape.top == 3096755
          and shape.width == 8048625):
        matriz_shapes.insert(0, shape)

if matriz_shapes:
    ref = matriz_shapes[0]
    L, T, W, H = ref.left, ref.top, ref.width, ref.height
    for s in matriz_shapes:
        s._element.getparent().remove(s._element)
    slide10.shapes.add_picture(MATRIZ, L, T, W, H)
    print(f"Slide 10 — matriz_confusao_int8.png adicionada ({W//914400}x{H//914400} in)")
else:
    print("AVISO: placeholder da matriz de confusao nao encontrado")

# ── 3. SLIDE 19 — screenshot do app + foto do hardware ──────────────────────
slide19 = prs.slides[18]

app_area = None
for shape in slide19.shapes:
    if (shape.has_text_frame
            and shape.text_frame.text.strip() == ""
            and shape.width > 5_000_000
            and shape.height > 5_000_000
            and shape.left < 2_000_000):
        app_area = shape
        break

if app_area:
    L, T, W, H = app_area.left, app_area.top, app_area.width, app_area.height
    app_area._element.getparent().remove(app_area._element)

    # Terço superior → app, terço inferior → hardware
    app_h = int(H * 0.62)
    hw_h  = H - app_h

    slide19.shapes.add_picture(APP_IOT,  L, T,          W, app_h)
    slide19.shapes.add_picture(HARDWARE, L, T + app_h,  W, hw_h)
    print("Slide 19 — app_iot.jpg + hardware_setup.jpg adicionados")
else:
    print("AVISO: área de imagem do slide 19 não encontrada")

# ── 4. Salvar PPTX intermediário ────────────────────────────────────────────
TMP = OUTPUT + ".tmp"
prs.save(TMP)
print("Imagens inseridas — salvando intermediário…")

# ── 5. SLIDE 20 — embutir vídeo via manipulação de ZIP/XML ──────────────────
#  Abre o PPTX como ZIP, injeta o mp4 e cria o elemento <p:pic> de vídeo.

NS = {
    "p":   "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a":   "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r":   "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "a14": "http://schemas.microsoft.com/office/drawing/2010/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
}

REL_VIDEO = "http://schemas.microsoft.com/office/2007/relationships/media"
CT_MP4    = "video/mp4"

# Coordenadas da área de vídeo no slide 20 (Shape 20)
VID_X, VID_Y, VID_CX, VID_CY = 10256490, 1429141, 7269510, 4095750

def find_slide_path(zf, slide_idx):
    """Retorna caminho interno do slide (ex: ppt/slides/slide20.xml)."""
    # lê ppt/presentation.xml para obter a ordem dos slides
    prs_xml = etree.fromstring(zf.read("ppt/presentation.xml"))
    sldIdLst = prs_xml.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst")
    sldId = list(sldIdLst)[slide_idx]
    rid = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    rels_xml = etree.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
    for rel in rels_xml:
        if rel.get("Id") == rid:
            target = rel.get("Target")
            return "ppt/" + target if not target.startswith("ppt/") else target
    raise ValueError(f"Slide {slide_idx} não encontrado")

def make_thumbnail_png():
    """Gera um PNG 640x360 preto (poster frame do vídeo) em memória."""
    import struct, zlib
    w, h = 640, 360
    raw = b"\x00" + b"\x00\x00\x00" * w  # filtro None + pixels RGB preto
    raw = raw * h
    def chunk(name, data):
        c = name + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
    png  = b"\x89PNG\r\n\x1a\n"
    png += chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
    png += chunk(b"IDAT", zlib.compress(raw))
    png += chunk(b"IEND", b"")
    return png

def inject_video(tmp_path, video_path, output_path):
    video_name   = os.path.basename(video_path)
    media_vid    = f"ppt/media/{video_name}"
    media_thumb  = "ppt/media/video_thumb.png"
    thumb_bytes  = make_thumbnail_png()

    with zipfile.ZipFile(tmp_path, "r") as zin:
        slide_path = find_slide_path(zin, 19)   # slide 20 = índice 19
        rels_path  = slide_path.replace("slides/", "slides/_rels/") + ".rels"
        slide_xml  = zin.read(slide_path)
        rels_xml   = zin.read(rels_path)
        ct_xml     = zin.read("[Content_Types].xml")

    # ── 5a. Relacionamentos: vídeo + thumbnail ─────────────────────────────
    rels_root = etree.fromstring(rels_xml)
    existing_ids = {r.get("Id") for r in rels_root}

    def next_rid(base):
        rid = base + "1"
        n = 1
        while rid in existing_ids:
            n += 1; rid = base + str(n)
        existing_ids.add(rid)
        return rid

    vid_rid   = next_rid("rIdVid")
    thumb_rid = next_rid("rIdThumb")

    etree.SubElement(rels_root, "Relationship", {
        "Id": vid_rid, "Type": REL_VIDEO,
        "Target": f"../media/{video_name}",
    })
    etree.SubElement(rels_root, "Relationship", {
        "Id": thumb_rid,
        "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
        "Target": "../media/video_thumb.png",
    })
    new_rels = etree.tostring(rels_root, xml_declaration=True,
                              encoding="UTF-8", standalone=True)

    # ── 5b. Content-Types ─────────────────────────────────────────────────
    ct_root = etree.fromstring(ct_xml)
    CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
    exts = {el.get("Extension") for el in ct_root.findall(f"{{{CT_NS}}}Default")}
    if "mp4" not in exts:
        etree.SubElement(ct_root, f"{{{CT_NS}}}Default",
                         {"Extension": "mp4", "ContentType": CT_MP4})
    if "png" not in exts:
        etree.SubElement(ct_root, f"{{{CT_NS}}}Default",
                         {"Extension": "png", "ContentType": "image/png"})
    new_ct = etree.tostring(ct_root, xml_declaration=True,
                             encoding="UTF-8", standalone=True)

    # ── 5c. Slide XML — substituir retângulo por <p:pic> de vídeo ─────────
    PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    RNS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    A14 = "http://schemas.microsoft.com/office/drawing/2010/main"
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"

    slide_root = etree.fromstring(slide_xml)
    sp_tree = slide_root.find(f".//{{{PNS}}}spTree")

    # Remover o retângulo placeholder (Shape 20) pelo offset X
    for sp in list(sp_tree.findall(f"{{{PNS}}}sp")):
        xfrm = sp.find(f".//{{{ANS}}}xfrm")
        if xfrm is not None:
            off = xfrm.find(f"{{{ANS}}}off")
            if off is not None and off.get("x") == str(VID_X):
                sp_tree.remove(sp)
                break

    VID_SHAPE_ID = "201"

    # ── <p:pic> de vídeo ──────────────────────────────────────────────────
    pic = etree.SubElement(sp_tree, f"{{{PNS}}}pic")

    nvPicPr = etree.SubElement(pic, f"{{{PNS}}}nvPicPr")

    # cNvPr — sem hlinkClick (timing cuida do play)
    cNvPr = etree.SubElement(nvPicPr, f"{{{PNS}}}cNvPr",
                              id=VID_SHAPE_ID, name="Video 1", descr=video_name)

    # cNvPicPr — referência a14:media (formato Office 2010+)
    cNvPicPr = etree.SubElement(nvPicPr, f"{{{PNS}}}cNvPicPr")
    etree.SubElement(cNvPicPr, f"{{{ANS}}}picLocks", noChangeAspect="1")
    cNvExtLst = etree.SubElement(cNvPicPr, f"{{{ANS}}}extLst")
    extA14 = etree.SubElement(cNvExtLst, f"{{{ANS}}}ext",
                               uri="{C0B2C2F2-D5CC-4A6E-B4CC-E8934CF15FC7}")
    a14m = etree.SubElement(extA14, f"{{{A14}}}media")
    a14m.set(f"{{{RNS}}}embed", vid_rid)

    # nvPr — SEM <p:ph> (standalone, não é placeholder de layout)
    #         apenas extLst com p14:media
    nvPr = etree.SubElement(nvPicPr, f"{{{PNS}}}nvPr")
    nvExtLst = etree.SubElement(nvPr, f"{{{PNS}}}extLst")
    extP14 = etree.SubElement(nvExtLst, f"{{{PNS}}}ext",
                               uri="{D42A27DB-BD31-4B8C-83A1-F6EECEF9847C}")
    p14m = etree.SubElement(extP14, f"{{{P14}}}media")
    p14m.set(f"{{{RNS}}}embed", vid_rid)

    # blipFill com thumbnail preto como poster frame
    blipFill = etree.SubElement(pic, f"{{{PNS}}}blipFill")
    blip = etree.SubElement(blipFill, f"{{{ANS}}}blip")
    blip.set(f"{{{RNS}}}embed", thumb_rid)
    stretch = etree.SubElement(blipFill, f"{{{ANS}}}stretch")
    etree.SubElement(stretch, f"{{{ANS}}}fillRect")

    spPr = etree.SubElement(pic, f"{{{PNS}}}spPr")
    xfrm = etree.SubElement(spPr, f"{{{ANS}}}xfrm")
    off  = etree.SubElement(xfrm, f"{{{ANS}}}off")
    off.set("x", str(VID_X)); off.set("y", str(VID_Y))
    ext2 = etree.SubElement(xfrm, f"{{{ANS}}}ext")
    ext2.set("cx", str(VID_CX)); ext2.set("cy", str(VID_CY))
    pg = etree.SubElement(spPr, f"{{{ANS}}}prstGeom", prst="rect")
    etree.SubElement(pg, f"{{{ANS}}}avLst")

    # ── <p:timing> — registra o vídeo como mídia click-to-play ───────────
    # Estrutura obrigatória para o PowerPoint reconhecer clique no vídeo.
    # spid deve bater com cNvPr id acima.
    timing_xml = f"""<p:timing xmlns:p="{PNS}" xmlns:a="{ANS}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" dur="indefinite" fill="hold">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:video>
                                    <p:cMediaNode vol="80000">
                                      <p:cTn id="6" dur="indefinite" fill="hold" display="0"/>
                                      <p:tgtEl><p:spTgt spid="{VID_SHAPE_ID}"/></p:tgtEl>
                                    </p:cMediaNode>
                                  </p:video>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrevClick" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNextClick" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst>
    <p:bldP spid="{VID_SHAPE_ID}" grpId="0"/>
  </p:bldLst>
</p:timing>"""

    timing_el = etree.fromstring(timing_xml)
    # Remover timing existente (se houver) e adicionar o novo como filho de <p:sld>
    sld_root = slide_root  # slide_root é <p:sld>
    for old_timing in sld_root.findall(f"{{{PNS}}}timing"):
        sld_root.remove(old_timing)
    sld_root.append(timing_el)

    new_slide = etree.tostring(slide_root, xml_declaration=True,
                               encoding="UTF-8", standalone=True)

    # ── 5d. Reempacotar ───────────────────────────────────────────────────
    with zipfile.ZipFile(tmp_path, "r") as zin, \
         zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            if item == slide_path:
                zout.writestr(item, new_slide)
            elif item == rels_path:
                zout.writestr(item, new_rels)
            elif item == "[Content_Types].xml":
                zout.writestr(item, new_ct)
            else:
                zout.writestr(item, zin.read(item))
        zout.write(video_path, media_vid)
        zout.writestr(media_thumb, thumb_bytes)

    print(f"Slide 20 — video embutido ({os.path.getsize(video_path)//1024//1024} MB) + thumbnail OK")

inject_video(TMP, VIDEO, OUTPUT)
os.remove(TMP)
print(f"\nPronto: {os.path.basename(OUTPUT)}")
